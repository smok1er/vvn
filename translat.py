import os

from docx.enum.style import WD_STYLE_TYPE
from docx.shared import Pt, RGBColor

import fitz
import telebot

from telebot import types
import requests

from pptx import Presentation
from docx import Document
import glob

print('tran')
r1 = requests.session()
idi = 209501902
token = '1215240244:AAGIv5jStQInzu7Qvx_XmPOwsEeyNT104fk'
bot = telebot.TeleBot("1215240244:AAGIv5jStQInzu7Qvx_XmPOwsEeyNT104fk")
open('tran.txt', 'a').write('')
c = []
g = []
d = []
work = []


@bot.message_handler(commands=['start'])
def key(msg):
    ch = msg.chat.id
    try:
        if ch in c:
            pass
        else:
            c.append(ch)

        if msg.chat.id == idi:
            admin = types.ReplyKeyboardMarkup()
            akoky = types.KeyboardButton('الاذاعة')
            zkoky = types.KeyboardButton('الاعضاء')
            admin.add(akoky, zkoky)
            bot.reply_to(msg, 'اهلا و سهلا بك ايها المطور {}'.format(msg.from_user.first_name), reply_markup=admin)

        else:
            idd = msg.chat.id
            res = \
                r1.get(f'https://api.telegram.org/bot{token}/getChatMember?chat_id=@akokybot&user_id={idd}').json()[
                    'result'][
                    'status']

            if res == 'left' and msg.chat.id not in [idi, 1490464385]:
                bot.send_message(ch, 'يجب عليك الاشتراك في القناة اول و من ثم استخدم البوت \nلطفا💚\n@akokybot')

            else:
                q = types.InlineKeyboardMarkup()

                a9 = types.InlineKeyboardButton("المطور 💚", url="https://t.me/Q5QQQQ")

                q.add(a9)

                bot.send_message(msg.chat.id, 'اهلا و سهلا بك في بالبوت الترجمة 💚', reply_markup=q)
            if f'{msg.chat.id}\n' in open('tran.txt', 'r'):
                pass
            else:
                open('tran.txt', 'a').write(f'{msg.chat.id}\n')
    except:
        pass


@bot.message_handler(content_types='text')
def an(msg):
    ch = msg.chat.id
    x = msg.text
    try:
        zzz = 's'
        if zzz == 'h':
            pass
        else:
            if f'{msg.chat.id}\n' in open('tran.txt', 'r'):
                pass
            else:
                open('tran.txt', 'a').write(f'{msg.chat.id}\n')
            idd = msg.chat.id
            res = \
                r1.get(f'https://api.telegram.org/bot{token}/getChatMember?chat_id=@akokybot&user_id={idd}').json()[
                    'result'][
                    'status']

            if res == 'left' and msg.chat.id not in [idi, 1490464385]:
                bot.send_message(ch, 'يجب عليك الاشتراك في القناة اول و من ثم استخدم البوت \nلطفا💚\n@akokybot')

            else:
                try:
                    if msg.text == 'الاعضاء' and msg.chat.id in [idi, 1490464385]:
                        x = open('tran.txt', 'r').readlines()
                        bot.reply_to(msg, '{}'.format(len(x)))
                except:
                    pass
                try:
                    if msg.text == 'الاذاعة' and msg.chat.id == idi:
                        markup = types.ForceReply(selective=False)
                        bot.send_message(msg.chat.id, "ارسل اذاعتك", reply_markup=markup)
                except:
                    pass
                try:
                    if msg.reply_to_message.text == "ارسل اذاعتك":
                        try:
                            s = msg.text
                            x = open('tran.txt', 'r')
                            for i in x:
                                try:
                                    bot.send_message(i.replace('\n', ''), s)
                                except:
                                    pass
                        except:
                            pass
                except:
                    pass
                w = msg.text.lower()
                try:
                    if msg.reply_to_message.text == "ارسل المصطلح الطبي المراد ترجمته":
                        if w[0] == 'ا' or w[0] == 'أ' and w[1] == 'ل':
                            url = f'https://context.reverso.net/translation/arabic-english/{w}'
                            head = {
                                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                'Accept-Encoding': 'gzip, deflate, br',
                                'Accept-Language': 'en-US,en;q=0.9',
                                'Connection': 'keep-alive',
                                'Cookie': 'didomi_token=eyJ1c2VyX2lkIjoiMTc1YWYxY2QtM2E4NS02NWNhLTlhNWYtNWMzMTMxYmY0OTQyIiwiY3JlYXRlZCI6IjIwMjAtMTEtMDlUMjI6MjU6MDMuODM3WiIsInVwZGF0ZWQiOiIyMDIwLTExLTA5VDIyOjI1OjAzLjgzN1oiLCJ2ZXJzaW9uIjpudWxsfQ==; __qca=P0-2069518181-1604960703577; _ga=GA1.2.1022455205.1604960704; _fbp=fb.1.1604960704785.1696161273; __gads=ID=6a415d82c09e7f2f:T=1604961014:S=ALNI_MYnjNMXo7nFRuuahoJIpqXamOQFvA; CTXTNODEID=bstweb12; experiment_context_N7gT3vKzX=0; JSESSIONID=kamXv4rTSPJG1mzG5D0fBsyY.bst-web12; context.lastpair=ar-en; history_entry=psycho]#[gastric]#[{w}; history_pair=en-ar]#[en-ar]#[ar-en; experiment_context_E3de3pqAZ=1; _gid=GA1.2.1607272187.1606403162; context.dapppromotion-count2=1; context.dapppromotion2=0; _gat_gtag_UA_2834324_41=1',
                                'Host': 'context.reverso.net',
                                'Referer': 'https://context.reverso.net/translation/arabic-english/%7Bw%7D',
                                'Sec-Fetch-Dest': 'document',
                                'Sec-Fetch-Mode': 'navigate',
                                'Sec-Fetch-Site': 'same-origin',
                                'Sec-Fetch-User': '?1',
                                'Upgrade-Insecure-Requests': '1',
                                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36'
                            }
                            j = r1.get(url=url, headers=head)
                            v = (j.text[
                                 j.text.find('<button class="other-content" data-other="0" data-negative="') + len(
                                     '<button class="other-content" data-other="0" data-negative="'):j.text.find(
                                     'Other translations</button>')])[:-2]
                            c = (v.replace("-{", "").replace("}", "\n"))
                            bot.send_message(msg.chat.id, f'* {c} *')

                        else:
                            try:
                                if f'{w[0]}' in ['a', 'b', 'c', 'd', 'e', 'f', 'g', 'h', 'i', 'j', 'k', 'l', 'm', 'n',
                                                 'o',
                                                 'p',
                                                 'q',
                                                 'r',
                                                 's', 't', 'u', 'v', 'w', 'x', 'y', 'z']:
                                    d = w
                                    v = f'https://www.tbeeb.net/%D9%82%D8%A7%D9%85%D9%88%D8%B3-%D8%B7%D8%A8%D9%8A/search.php?q={d}+&dictionary=%D8%A8%D8%AD%D8%AB'
                                    try:
                                        i = 1
                                        zz = ''
                                        while i < len(r1.get(v).text.split('"tden">')):
                                            if i > 7:
                                                pass
                                            else:
                                                x1 = r1.get(v).text.split('"tden">')[i][
                                                     :r1.get(v).text.split('"tden">')[i].find('</td>')]
                                                x2 = r1.get(v).text.split('"tdar">')[i][
                                                     :r1.get(v).text.split('"tdar">')[i].find('</td>')]
                                                zz += f'{x2}\n{x1}\n'
                                            i += 1

                                        bot.send_message(ch, zz)
                                    except:
                                        try:
                                            url = f'https://context.reverso.net/translation/arabic-english/{w}'
                                            head = {
                                                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                                'Accept-Encoding': 'gzip, deflate, br',
                                                'Accept-Language': 'en-US,en;q=0.9',
                                                'Connection': 'keep-alive',
                                                'Cookie': 'didomi_token=eyJ1c2VyX2lkIjoiMTc1YWYxY2QtM2E4NS02NWNhLTlhNWYtNWMzMTMxYmY0OTQyIiwiY3JlYXRlZCI6IjIwMjAtMTEtMDlUMjI6MjU6MDMuODM3WiIsInVwZGF0ZWQiOiIyMDIwLTExLTA5VDIyOjI1OjAzLjgzN1oiLCJ2ZXJzaW9uIjpudWxsfQ==; __qca=P0-2069518181-1604960703577; _ga=GA1.2.1022455205.1604960704; _fbp=fb.1.1604960704785.1696161273; __gads=ID=6a415d82c09e7f2f:T=1604961014:S=ALNI_MYnjNMXo7nFRuuahoJIpqXamOQFvA; CTXTNODEID=bstweb12; experiment_context_N7gT3vKzX=0; JSESSIONID=kamXv4rTSPJG1mzG5D0fBsyY.bst-web12; context.lastpair=ar-en; history_entry=psycho]#[gastric]#[{w}; history_pair=en-ar]#[en-ar]#[ar-en; experiment_context_E3de3pqAZ=1; _gid=GA1.2.1607272187.1606403162; context.dapppromotion-count2=1; context.dapppromotion2=0; _gat_gtag_UA_2834324_41=1',
                                                'Host': 'context.reverso.net',
                                                'Referer': 'https://context.reverso.net/translation/arabic-english/%7Bw%7D',
                                                'Sec-Fetch-Dest': 'document',
                                                'Sec-Fetch-Mode': 'navigate',
                                                'Sec-Fetch-Site': 'same-origin',
                                                'Sec-Fetch-User': '?1',
                                                'Upgrade-Insecure-Requests': '1',
                                                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36'
                                            }
                                            j = r1.get(url=url, headers=head)
                                            v = (j.text[
                                                 j.text.find(
                                                     '<button class="other-content" data-other="0" data-negative="') + len(
                                                     '<button class="other-content" data-other="0" data-negative="'):j.text.find(
                                                     'Other translations</button>')])[:-2]
                                            c = (v.replace("-{", "").replace("}", "\n"))
                                            bot.send_message(msg.chat.id, f'* {c} *')
                                        except:
                                            pass
                                    head = {
                                        'accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                        'accept-encoding': 'gzip, deflate, br',
                                        'accept-language': 'en-US,en;q=0.9',
                                        'cookie': 'CGIC=IocBdGV4dC9odG1sLGFwcGxpY2F0aW9uL3hodG1sK3htbCxhcHBsaWNhdGlvbi94bWw7cT0wLjksaW1hZ2UvYXZpZixpbWFnZS93ZWJwLGltYWdlL2FwbmcsKi8qO3E9MC44LGFwcGxpY2F0aW9uL3NpZ25lZC1leGNoYW5nZTt2PWIzO3E9MC45; HSID=ACpCi--xKNWBuzc1V; SSID=AxTv5nnKkxKd-t24h; APISID=OvtsQIlYX38DaDAB/AejTrmzLNAmxJKW5L; SAPISID=VictV1NMEEknSmKC/Ap3tCyc9Rw-nNhkKN; __Secure-3PAPISID=VictV1NMEEknSmKC/Ap3tCyc9Rw-nNhkKN; CONSENT=YES+IQ.ar+201908; SEARCH_SAMESITE=CgQIoJAB; OTZ=5713657_44_44__44_; SID=3gekKKpt2BDYs3HGt2dbGMJygRXv0S_uH30dI87nwwd9kPVGeJPLKTJluqScaUfUXXkqGg.; __Secure-3PSID=3gekKKpt2BDYs3HGt2dbGMJygRXv0S_uH30dI87nwwd9kPVGofFKYNMh0FN74vl_7RXMLg.; ANID=AHWqTUl-D-kiT0ekElpQYvwaLkYpmfjEj0_vWVNyeQ9lR3vPGLYG1ea_vAtrh9jw; 1P_JAR=2020-11-26-23; NID=204=Nvh8sYQiAGl2xGl9Rivw7mwHK0oiirbBuBjelwbKer7spRkQS3Xh-9GYO5GR2WOwUG-GICaT4wz4vP8fWCpH8Zao3AEmb_XGvZ7CoLvMiTCbcSnMjY6GCGnFcwc9Ap2D0Iu1ltcyj4qzLl25BPgZf0KdrrpVYB_UyBb_LrJWA7A4dgrtBLEwtyBFzHe1XdYKS_yOA6QzeeHlNNtZm2YYpYUzvzevK2wxD2EMs9f9OOnq_es; DV=Q77RAgvgUI1KQP_f7OZE3e4dEcltYJdw4P5MyRCtgwEAAHAXk3Gqw6shngAAAAyZb21iQXUwRQAAAA; SIDCC=AJi4QfGRLCnLLJayzTS67kgRZvvsMJ8WJx4Z3VQIypirxBw8XSvY82_9ydx6j9Wr0SSSgWwkPw; __Secure-3PSIDCC=AJi4QfFg7EgRAZ_f0vxbunVqGgMst8xwoewWB6J0ZP6y8sihJ3OufFh34ZiNnfFg87jxfwhwgg',
                                        'referer': 'https://www.google.com/',
                                        'sec-fetch-dest': 'document',
                                        'sec-fetch-mode': 'navigate',
                                        'sec-fetch-site': 'same-origin',
                                        'sec-fetch-user': '?1',
                                        'upgrade-insecure-requests': '1',
                                        'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36',
                                        'x-client-data': 'CIy2yQEIpLbJAQjBtskBCKmdygEIlqzKAQisx8oBCPbHygEI6cjKAQi0y8oBCI3PygEI3NXKAQjul8sBCJGZywEImJrLARiKwcoB'
                                    }
                                    r = requests.get(
                                        f'https://www.google.com/search?q={d}+anatomy&source=lnms&tbm=isch&sa=X&ved=2ahUKEwjfv9Xg5YrvAhWHnxQKHe-2AfEQ_AUoAXoECBQQAw&biw=1821&bih=876#imgrc=3PIl4efjvoyhfM',
                                        headers=head)
                                    i = 0
                                    import shutil
                                    try:
                                        while i < 3:
                                            try:
                                                x = r.text.split('jpg"')[i].split('["http')[-1]
                                                v = f'http{x}jpg'

                                                with open(fr'koky{i}.jpg', 'wb') as mp3:
                                                    shutil.copyfileobj(r1.get(v, stream=True).raw, mp3)
                                                    bot.send_photo(chat_id=ch,
                                                                   photo=open(f'koky{i}.jpg', 'rb'))

                                            except:
                                                pass
                                            i += 1
                                    except:
                                        pass


                                else:
                                    url = f'https://context.reverso.net/translation/arabic-english/{w}'
                                    head = {
                                        'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                        'Accept-Encoding': 'gzip, deflate, br',
                                        'Accept-Language': 'en-US,en;q=0.9',
                                        'Connection': 'keep-alive',
                                        'Cookie': 'didomi_token=eyJ1c2VyX2lkIjoiMTc1YWYxY2QtM2E4NS02NWNhLTlhNWYtNWMzMTMxYmY0OTQyIiwiY3JlYXRlZCI6IjIwMjAtMTEtMDlUMjI6MjU6MDMuODM3WiIsInVwZGF0ZWQiOiIyMDIwLTExLTA5VDIyOjI1OjAzLjgzN1oiLCJ2ZXJzaW9uIjpudWxsfQ==; __qca=P0-2069518181-1604960703577; _ga=GA1.2.1022455205.1604960704; _fbp=fb.1.1604960704785.1696161273; __gads=ID=6a415d82c09e7f2f:T=1604961014:S=ALNI_MYnjNMXo7nFRuuahoJIpqXamOQFvA; CTXTNODEID=bstweb12; experiment_context_N7gT3vKzX=0; JSESSIONID=kamXv4rTSPJG1mzG5D0fBsyY.bst-web12; context.lastpair=ar-en; history_entry=psycho]#[gastric]#[{w}; history_pair=en-ar]#[en-ar]#[ar-en; experiment_context_E3de3pqAZ=1; _gid=GA1.2.1607272187.1606403162; context.dapppromotion-count2=1; context.dapppromotion2=0; _gat_gtag_UA_2834324_41=1',
                                        'Host': 'context.reverso.net',
                                        'Referer': 'https://context.reverso.net/translation/arabic-english/%7Bw%7D',
                                        'Sec-Fetch-Dest': 'document',
                                        'Sec-Fetch-Mode': 'navigate',
                                        'Sec-Fetch-Site': 'same-origin',
                                        'Sec-Fetch-User': '?1',
                                        'Upgrade-Insecure-Requests': '1',
                                        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36'
                                    }
                                    j = r1.get(url=url, headers=head)
                                    v = (j.text[
                                         j.text.find(
                                             '<button class="other-content" data-other="0" data-negative="') + len(
                                             '<button class="other-content" data-other="0" data-negative="'):j.text.find(
                                             'Other translations</button>')])[:-2]
                                    c = (v.replace("-{", "").replace("}", "\n"))
                                    bot.send_message(msg.chat.id, f'* {c} *')
                            except:
                                pass


                except:
                    pass
                if msg.reply_to_message:
                    pass
                else:
                    try:
                        url = 'https://www.arabtran.com/gtranslate/'

                        head = {
                            'accept': '*/*',
                            'accept-encoding': 'gzip, deflate, br',
                            'accept-language': 'en-US,en;q=0.9',
                            'content-length': '31',
                            'content-type': 'application/x-www-form-urlencoded; charset=UTF-8',
                            'cookie': '_ga=GA1.2.189101859.1607603223; _gid=GA1.2.154684385.1607603223; __gads=ID=ba2b5a0fe52a4c1b-22c3694088a60095:T=1607603223:RT=1607603223:S=ALNI_MbdZ0H2MeiuATaH2_vAl4hrUnPm8Q; _gat=1; sc_is_visitor_unique=rx12068393.1607603329.CA2EB2FDB25A4FB19C08D5DA0B092EBC.1.1.1.1.1.1.1.1.1',
                            'origin': 'https://www.arabtran.com',
                            'referer': 'https://www.arabtran.com/tarjamat_anjilizi_earabiun/',
                            'sec-fetch-dest': 'empty',
                            'sec-fetch-mode': 'cors',
                            'sec-fetch-site': 'same-origin',
                            'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/87.0.4280.88 Safari/537.36',
                            'x-requested-with': 'XMLHttpRequest',
                        }

                        if f'{w[0]}' in ['a', 'b', 'c', 'd', 'e', 'f', 'g', 'h', 'i', 'j', 'k', 'l', 'm', 'n', 'o', 'p',
                                         'q',
                                         'r',
                                         's', 't', 'u', 'v', 'w', 'x', 'y', 'z'] or f'{w[1]}' in ['a', 'b', 'c', 'd',
                                                                                                  'e',
                                                                                                  'f',
                                                                                                  'g',
                                                                                                  'h', 'i', 'j', 'k',
                                                                                                  'l',
                                                                                                  'm',
                                                                                                  'n',
                                                                                                  'o', 'p', 'q', 'r',
                                                                                                  's',
                                                                                                  't',
                                                                                                  'u',
                                                                                                  'v', 'w', 'x', 'y',
                                                                                                  'z'] or f'{w[2]}' in [
                            'a',
                            'b',
                            'c',
                            'd',
                            'e',
                            'f',
                            'g',
                            'h',
                            'i',
                            'j',
                            'k',
                            'l',
                            'm',
                            'n',
                            'o',
                            'p',
                            'q',
                            'r',
                            's',
                            't',
                            'u',
                            'v',
                            'w',
                            'x',
                            'y',
                            'z']:
                            data = {
                                'text': w,
                                'gfrom': 'en',
                                'gto': 'ar',
                                'key': 'ABC'
                            }
                            x = r1.post(url=url, data=data, ).text
                            bot.send_message(msg.chat.id, x)

                        else:
                            data = {
                                'text': w,
                                'gfrom': 'ar',
                                'gto': 'en',
                                'key': 'ABC'
                            }
                            x = r1.post(url=url, data=data, ).text
                            bot.send_message(msg.chat.id, x)

                    except:
                        pass
    except:
        pass


@bot.callback_query_handler(lambda call: True)
def any(call):
    if call.data == 'q1':
        markup = types.ForceReply(selective=False)
        bot.send_message(call.message.chat.id, "ارسل المصطلح الطبي المراد ترجمته", reply_markup=markup)


h = []


@bot.message_handler(content_types='document')
def an(msg):
    ch = msg.chat.id  #
    try:
        if ch in work:
            bot.send_message(ch, 'انتضر قليلا ليترجم الملف الذي ارسلته سابقا⚠️')
        else:
            work.append(ch)

            if ch in c:
                pass
            else:
                c.append(ch)

            if f'{msg.chat.id}\n' in open('tran.txt', 'r'):
                pass
            else:
                open('tran.txt', 'a').write(f'{msg.chat.id}\n')
            idd = msg.chat.id
            res = \
                r1.get(f'https://api.telegram.org/bot{token}/getChatMember?chat_id=@akokybot&user_id={idd}').json()[
                    'result'][
                    'status']

            if res == 'left' and msg.chat.id not in [idi, 1490464385]:
                bot.send_message(ch, 'يجب عليك الاشتراك في القناة اول و من ثم استخدم البوت \nلطفا💚\n@akokybot')

            else:
                ch = msg.chat.id
                user = msg.from_user.username
                open(f'{user}.txt', 'w').write('')
                open(f'{user}ar.txt', 'w', encoding='utf-8').write('')

                try:
                    file_info = bot.get_file(msg.document.file_id)
                    downloaded_file = bot.download_file(file_info.file_path)
                    x = msg.document.file_name[-4:]
                    if x == '.pdf' or x == '.PDF':

                        if ch in h:
                            bot.send_message(ch, 'انتضر قليلا من فضلك لكي يتم التحميل 💚')
                        else:
                            h.append(ch)
                            bot.send_message(ch, 'حسنا, سيتم التحميل الرجاء انتظار دقيقة واحدة 💚')
                            with open('ahmed.pdf', 'wb') as new_file:
                                new_file.write(downloaded_file)

                            with fitz.open("ahmed.pdf") as doc:
                                text = ""
                                open(f'{user}.txt', 'a').write(
                                    'Done translation by bot ahmed \n@eng_ahmed8_bot\n@Q5QQQQ\n')
                                open(f'{user}.txt', 'a').write('-----------------------\n')

                                for page in doc:
                                    text += page.getText()
                                    text += f'----------------{str(page)[:7]}-------------------\n'
                                    text += '|#|'
                                    # open(f'{user}.txt', 'a').write(f'{text}\n')
                                open(f'{user}.txt', 'a', encoding='utf-8').write(text)

                                z = open(f'{user}.txt', 'r', encoding='utf-8').read()
                                url = 'https://www.arabtran.com/gtranslate/'
                                head = {
                                    'accept': '*/*',
                                    'accept-encoding': 'gzip, deflate, br',
                                    'accept-language': 'en-US,en;q=0.9',
                                    'content-length': '31',
                                    'content-type': 'application/x-www-form-urlencoded; charset=UTF-8',
                                    'cookie': '_ga=GA1.2.189101859.1607603223; _gid=GA1.2.154684385.1607603223; __gads=ID=ba2b5a0fe52a4c1b-22c3694088a60095:T=1607603223:RT=1607603223:S=ALNI_MbdZ0H2MeiuATaH2_vAl4hrUnPm8Q; _gat=1; sc_is_visitor_unique=rx12068393.1607603329.CA2EB2FDB25A4FB19C08D5DA0B092EBC.1.1.1.1.1.1.1.1.1',
                                    'origin': 'https://www.arabtran.com',
                                    'referer': 'https://www.arabtran.com/tarjamat_anjilizi_earabiun/',
                                    'sec-fetch-dest': 'empty',
                                    'sec-fetch-mode': 'cors',
                                    'sec-fetch-site': 'same-origin',
                                    'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/87.0.4280.88 Safari/537.36',
                                    'x-requested-with': 'XMLHttpRequest',
                                }
                                z = open(f'{user}.txt', 'r', encoding='utf-8').read()

                                if len(z) > 5000:
                                    if ch in d:
                                        k = d.index(ch)
                                        try:
                                            d.remove(d[k])
                                            d.remove(d[k])
                                        except:
                                            pass
                                        d.append(ch)
                                        d.append(' ')
                                    else:
                                        d.append(ch)
                                        d.append(' ')

                                    t = ''
                                    e = d[d.index(ch) + 1]
                                    u = []
                                    p = 0
                                    while (p - 1) * 5000 < len(z):
                                        if p == 0:

                                            data = {
                                                'text': z,
                                                'gfrom': 'en',
                                                'gto': 'ar',
                                                'key': 'ABC'
                                            }
                                            x = r1.post(url=url, data=data, ).text
                                            t += x
                                            q = t.split('\n')
                                            q1 = z.split('\n')
                                            i = 0
                                            while i < len(q1):
                                                try:
                                                    open(f'{user}ar.txt', 'a', encoding='utf-8').write(
                                                        f'{q[i]}\n')
                                                    e += f'{q1[i]}\n'
                                                    if i - len(q1) == -1:
                                                        d.insert(d.index(ch) + 1, e)
                                                except:
                                                    pass
                                                i += 1
                                        else:

                                            data = {
                                                'text': z[z.find(d[d.index(ch) + 1]) + len(e):],
                                                'gfrom': 'en',
                                                'gto': 'ar',
                                                'key': 'ABC'
                                            }
                                            x = r1.post(url=url, data=data, ).text
                                            q = x.split('\n')
                                            q1 = z[z.find(e) + len(e):].split('\n')
                                            i = 0
                                            while i < len(q1):
                                                try:
                                                    open(f'{user}ar.txt', 'a', encoding='utf-8').write(
                                                        f'{q[i]}\n')
                                                    e += f'{q1[i]}\n'
                                                    if i - len(q1) == -1:
                                                        d.insert(d.index(ch) + 1, e)
                                                except:
                                                    pass
                                                i += 1
                                        p += 1
                                else:
                                    data = {
                                        'text': z,
                                        'gfrom': 'en',
                                        'gto': 'ar',
                                        'key': 'ABC'
                                    }
                                    x = r1.post(url=url, data=data, ).text.split('\n')
                                    xx = z.split('\n')
                                    i = 0
                                    while i < len(x):
                                        open(f'{user}ar.txt', 'a', encoding='utf-8').write(f"{x[i]}\n")
                                        i += 1
                                x1 = open(f'{user}ar.txt', 'r', encoding='utf-8').read()
                                open('ahmed.txt', 'w', encoding='utf-8').write(x1)
                                bot.send_document(ch, open(f'ahmed.txt', 'r', encoding='utf-8'),
                                                  caption=f'{msg.document.file_name}\nترجمة عربي فقط')
                                x2 = open(f'{user}.txt', 'r', encoding='utf-8').read()
                                cc = 0
                                u = ''
                                while cc < len(x2.split('|#|')):
                                    try:
                                        u += f"{z.split('|#|')[cc]}\n{x1.split('| # |')[cc]}\n"
                                    except:
                                        pass
                                    cc += 1
                                document = Document()
                                paragraph = document.add_paragraph(str(u))
                                paragraph.style = document.styles.add_style('', WD_STYLE_TYPE.PARAGRAPH)
                                font = paragraph.style.font
                                font.size = Pt(16)
                                font.bold = True
                                font.color.rgb = RGBColor(0, 0, 0)
                                document.save('ahmed1.docx')
                                bot.send_document(ch, open(f'ahmed1.docx', 'rb'),
                                                  caption=f'{msg.document.file_name}\nترجمة صفحة صفحة')
                                cc = 0
                                u = ''
                                open(f'ahmed1.txt', 'w', encoding='utf-8').write(f"")
                                x1 = open(f'{user}ar.txt', 'r', encoding='utf-8').read()
                                z = open(f'{user}.txt', 'r', encoding='utf-8').read()
                                while cc < len(x1.split('\n')):
                                    try:
                                        if z.split('\n')[cc] == '':
                                            pass
                                        elif len(z.split('\n')[cc].split(' ')) < 4:
                                            x11 = x1.split('\n')[cc]
                                            xz = z.split('\n')[cc]
                                            try:
                                                u += f"{x11}\n{xz}\n"
                                                open(f'ahmed1.txt', 'a', encoding='utf-8').write(f"{x11} {xz}\n")
                                            except:
                                                pass
                                        else:
                                            x11 = x1.split('\n')[cc]
                                            xz = z.split('\n')[cc]
                                            try:
                                                u += f"{x11}\n{xz}\n"
                                                open(f'ahmed1.txt', 'a', encoding='utf-8').write(f"{x11}\n{xz}\n")
                                            except:
                                                pass
                                    except:
                                        pass
                                    cc += 1
                                document1 = Document()
                                paragraph = document1.add_paragraph(
                                    str(open(f'ahmed1.txt', 'r', encoding='utf-8').read()))
                                paragraph.style = document1.styles.add_style('', WD_STYLE_TYPE.PARAGRAPH)
                                font = paragraph.style.font
                                font.size = Pt(16)
                                font.bold = True
                                font.color.rgb = RGBColor(0, 0, 0)
                                document1.save('ahmed2.docx')
                                bot.send_document(ch, open(f'ahmed2.docx', 'rb'),
                                                  caption=f'{msg.document.file_name}\nترجمة سطرية')


                    else:
                        pp = []
                        open(f'{user}.txt', 'w').write('')
                        x = (msg.document.file_name).split('.')[-1]
                        if x == 'pptx' or x == 'PPTX':
                            with open('ahmed.pptx', 'wb') as new_file:
                                new_file.write(downloaded_file)
                            for eachfile in glob.glob("*.pptx"):
                                prs = Presentation(eachfile)
                                open(f'{user}.txt', 'a').write(
                                    'Done translation by ahmed bot \n@eng_ahmed8_bot\n@Q5QQQQ\n')
                                open(f'{user}.txt', 'a').write('-----------------------\n')
                                u = ''
                                for slide in prs.slides:

                                    for shape in slide.shapes:
                                        if hasattr(shape, "text"):
                                            u += shape.text

                                pp.append(u)
                        elif x == 'docx' or x == 'DOCX':
                            u = ''
                            with open('ahmed.docx', 'wb') as new_file:
                                new_file.write(downloaded_file)
                            document = Document('ahmed.docx')
                            open(f'{user}.txt', 'a', encoding='utf-8').write(
                                'Done translation by ahmed bot \n@eng_ahmed8_bot\n@Q5QQQQ\n')
                            open(f'{user}.txt', 'a', encoding='utf-8').write('-----------------------\n')
                            for p in document.paragraphs:
                                u += p.text
                                u += '\n'
                            pp.append(u)

                        else:
                            bot.send_message(ch, 'عذرا لم يتم ترجمة ملفك لانه ليس✳️ \npdf or pptx or docx\n')

                        open(f'{user}.txt', 'a', encoding='utf-8').write(pp[0])
                        url = 'https://www.arabtran.com/gtranslate/'
                        head = {
                            'accept': '*/*',
                            'accept-encoding': 'gzip, deflate, br',
                            'accept-language': 'en-US,en;q=0.9',
                            'content-length': '31',
                            'content-type': 'application/x-www-form-urlencoded; charset=UTF-8',
                            'cookie': '_ga=GA1.2.189101859.1607603223; _gid=GA1.2.154684385.1607603223; __gads=ID=ba2b5a0fe52a4c1b-22c3694088a60095:T=1607603223:RT=1607603223:S=ALNI_MbdZ0H2MeiuATaH2_vAl4hrUnPm8Q; _gat=1; sc_is_visitor_unique=rx12068393.1607603329.CA2EB2FDB25A4FB19C08D5DA0B092EBC.1.1.1.1.1.1.1.1.1',
                            'origin': 'https://www.arabtran.com',
                            'referer': 'https://www.arabtran.com/tarjamat_anjilizi_earabiun/',
                            'sec-fetch-dest': 'empty',
                            'sec-fetch-mode': 'cors',
                            'sec-fetch-site': 'same-origin',
                            'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/87.0.4280.88 Safari/537.36',
                            'x-requested-with': 'XMLHttpRequest',
                        }
                        z = open(f'{user}.txt', 'r', encoding='utf-8').read()

                        if len(z) > 5000:

                            if ch in d:
                                k = d.index(ch)
                                try:
                                    d.remove(d[k])
                                    d.remove(d[k])
                                except:
                                    pass
                                d.append(ch)
                                d.append(' ')
                            else:
                                d.append(ch)
                                d.append(' ')

                            t = ''
                            e = d[d.index(ch) + 1]
                            u = []
                            p = 0

                            while (p - 1) * 5000 < len(z):

                                if p == 0:

                                    data = {
                                        'text': z,
                                        'gfrom': 'en',
                                        'gto': 'ar',
                                        'key': 'ABC'
                                    }
                                    x = r1.post(url=url, data=data, ).text
                                    t += x
                                    q = t.split('\n')
                                    q1 = z.split('\n')
                                    i = 0
                                    while i < len(q1):
                                        try:
                                            open(f'{user}ar.txt', 'a', encoding='utf-8').write(
                                                f'{q[i]}\n{q1[i]}\n')
                                            e += f'{q1[i]}\n'
                                            if i - len(q1) == -1:
                                                d.insert(d.index(ch) + 1, e)
                                        except:
                                            pass
                                        i += 1
                                else:
                                    data = {
                                        'text': z[z.find(d[d.index(ch) + 1]) + len(e):],
                                        'gfrom': 'en',
                                        'gto': 'ar',
                                        'key': 'ABC'
                                    }
                                    x = r1.post(url=url, data=data, ).text
                                    q = x.split('\n')
                                    q1 = z[z.find(e) + len(e):].split('\n')
                                    i = 0
                                    while i < len(q1):
                                        try:
                                            open(f'{user}ar.txt', 'a', encoding='utf-8').write(
                                                f'{q[i]}\n{q1[i]}\n')
                                            e += f'{q1[i]}\n'
                                            if i - len(q1) == -1:
                                                d.insert(d.index(ch) + 1, e)
                                        except:
                                            pass
                                        i += 1
                                p += 1

                        else:
                            data = {
                                'text': z,
                                'gfrom': 'en',
                                'gto': 'ar',
                                'key': 'ABC'
                            }
                            x = r1.post(url=url, data=data, ).text.split('\n')
                            xx = z.split('\n')
                            i = 0
                            while i < len(x):
                                open(f'{user}ar.txt', 'a', encoding='utf-8').write(f"{x[i]}\n{xx[i]}\n")
                                i += 1
                        open('ahmed.txt', 'w', encoding='utf-8').write(
                            open(f'{user}ar.txt', 'r', encoding='utf-8').read())
                        bot.send_document(ch, open(f'ahmed.txt', 'r', encoding='utf-8'),
                                          caption=msg.document.file_name)
                        document = Document()

                        paragraph = document.add_paragraph(str(open(f'ahmed.txt', 'r', encoding='utf-8').read()))

                        paragraph.style = document.styles.add_style('', WD_STYLE_TYPE.PARAGRAPH)
                        font = paragraph.style.font
                        font.size = Pt(15)
                        font.bold = True
                        font.color.rgb = RGBColor(0, 0, 100)
                        document.save('ahmed1.docx')
                        bot.send_document(ch, open(f'ahmed1.docx', 'rb'),
                                          caption=msg.document.file_name)

                    try:
                        h.remove(ch)
                    except:
                        pass
                except:
                    pass

            q = types.InlineKeyboardMarkup()
            a9 = types.InlineKeyboardButton("المطور 💚", url="https://t.me/Q5QQQQ")

            q.add(a9)

            try:
                os.remove(f'{user}.txt')
                os.remove(f'{user}ar.txt')
            except:
                pass
            bot.send_message(ch, 'الإعلانات ⬇,️', reply_markup=q)
        try:
            work.remove(ch)
        except:
            pass
    except:
        pass


i = 0
while i == 0:
    try:
        bot.polling(none_stop=True)
        i += 1
    except:
        pass
